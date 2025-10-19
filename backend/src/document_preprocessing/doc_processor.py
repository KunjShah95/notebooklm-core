import os
import logging
from typing import List, Dict, Any, Optional
from dataclasses import dataclass
from pathlib import Path
import hashlib
from datetime import datetime

import pymupdf
from docx import Document as DocxDocument
from pptx import Presentation

logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

@dataclass
class DocumentMetadata:
    content: str
    source_file: str
    source_type: str  # 'pdf', 'txt', 'md', 'web', 'audio'
    page_number: Optional[int] = None
    chunk_index: int = 0
    start_char: Optional[int] = None
    end_char: Optional[int] = None
    metadata: Dict[str, Any] = None
    chunk_id: str = ""

    def __post_init__(self):
        if not self.chunk_id:
            self.chunk_id = self._generate_chunk_id()
        if self.metadata is None:
            self.metadata = {}

    def _generate_chunk_id(self) -> str:
        content_hash = hashlib.md5(self.content.encode()).hexdigest()[:8]
        return f"{self.source_type}_{self.chunk_index}_{content_hash}"

    def get_citation_info(self) -> Dict[str, Any]:
        citation = {
            'source': self.source_file,
            'type': self.source_type,
            'chunk_id': self.chunk_id,
            'chunk_index': self.chunk_index
        }

        if self.page_number:
            citation['page'] = self.page_number
        if self.start_char is not None or self.end_char is not None:
            citation['char_range'] = f"{self.start_char}-{self.end_char}"

        citation.update(self.metadata)
        return citation


# Backwards-compatible alias: some modules expect DocumentChunk
# which is equivalent to DocumentMetadata in this project.
DocumentChunk = DocumentMetadata

__all__ = [
    'DocumentMetadata',
    'DocumentChunk',
    'DocumentProcessor'
]

class DocumentProcessor:
    def __init__(self, chunk_size: int = 1000, chunk_overlap: int = 200):
        self.chunk_size = chunk_size
        self.chunk_overlap = chunk_overlap
        self.supported_formats = {'.pdf', '.txt', '.md', '.pptx', '.docx'}  # add other formats if need be

    def process_document(self, file_path: str) -> List[DocumentMetadata]:
        file_path = Path(file_path)

        if not file_path.exists():
            raise FileNotFoundError(f"File not found: {file_path}")
        if file_path.suffix.lower() not in self.supported_formats:
            raise ValueError(f"Unsupported file format: {file_path.suffix}")

        logger.info(f"Processing document: {file_path.name}")

        try:
            if file_path.suffix.lower() == '.pdf':
                return self._process_pdf(file_path)
            elif file_path.suffix.lower() in {'.txt', '.md'}:
                return self._process_text_file(file_path)
            elif file_path.suffix.lower() == '.pptx':
                return self._process_pptx(file_path)
            elif file_path.suffix.lower() == '.docx':
                return self._process_docx(file_path)

        except Exception as e:
            logger.error(f"Error processing document {file_path.name}: {e}")
            raise

    def _process_pdf(self, file_path: Path) -> List[DocumentMetadata]:
        chunks: List[DocumentMetadata] = []
        try:
            doc = pymupdf.open(str(file_path))
            total_pages = len(doc)

            for page_num in range(total_pages):
                page = doc.load_page(page_num)
                text = page.get_text()

                if not text.strip():
                    continue

                # Get page metadata
                page_metadata = {
                    'total_pages': total_pages,
                    'page_width': page.rect.width,
                    'page_height': page.rect.height,
                    'processed_at': datetime.now().isoformat()
                }

                page_chunks = self._create_chunks_from_text(
                    text,
                    file_path.name,
                    source_type='pdf',
                    page_number=page_num + 1,
                    additional_metadata=page_metadata
                )
                chunks.extend(page_chunks)

            doc.close()
            logger.info(f"Processed PDF: {len(chunks)} chunks from {total_pages} pages")

        except Exception as e:
            logger.error(f"Error processing PDF {file_path}: {str(e)}")
            raise

        return chunks

    def _process_text_file(self, file_path: Path) -> List[DocumentMetadata]:
        try:
            with open(file_path, 'r', encoding='utf-8') as file:
                content = file.read()

            metadata = {
                'file_size': file_path.stat().st_size,
                'encoding': 'utf-8',
                'processed_at': datetime.now().isoformat()
            }

            chunks = self._create_chunks_from_text(
                content,
                file_path.name,
                source_type=file_path.suffix.lower()[1:],  # 'txt' or 'md'
                page_number=None,
                additional_metadata=metadata
            )

            logger.info(f"Processed text file: {len(chunks)} chunks")
            return chunks

        except Exception as e:
            logger.error(f"Error processing text file {file_path}: {str(e)}")
            raise

    def _process_pptx(self, file_path: Path) -> List[DocumentMetadata]:
        try:
            presentation = Presentation(str(file_path))
            text_content = []

            for slide_number, slide in enumerate(presentation.slides, 1):
                slide_text = []
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        slide_text.append(shape.text)
                if slide_text:
                    text_content.append(f"Slide {slide_number}:\n" + "\n".join(slide_text))

            full_text = "\n\n".join(text_content)

            metadata = {
                'total_slides': len(presentation.slides),
                'file_size': file_path.stat().st_size,
                'processed_at': datetime.now().isoformat()
            }

            chunks = self._create_chunks_from_text(
                full_text,
                file_path.name,
                source_type='pptx',
                page_number=None,
                additional_metadata=metadata
            )

            logger.info(f"Processed PPTX: {len(chunks)} chunks from {len(presentation.slides)} slides")
            return chunks

        except Exception as e:
            logger.error(f"Error processing PPTX {file_path}: {str(e)}")
            raise

    def _process_docx(self, file_path: Path) -> List[DocumentMetadata]:
        try:
            doc = DocxDocument(str(file_path))
            text_content = []

            for para in doc.paragraphs:
                if para.text.strip():
                    text_content.append(para.text)

            # Also extract from tables
            for table in doc.tables:
                for row in table.rows:
                    for cell in row.cells:
                        if cell.text.strip():
                            text_content.append(cell.text)

            full_text = "\n".join(text_content)

            metadata = {
                'total_paragraphs': len(doc.paragraphs),
                'total_tables': len(doc.tables),
                'file_size': file_path.stat().st_size,
                'processed_at': datetime.now().isoformat()
            }

            chunks = self._create_chunks_from_text(
                full_text,
                file_path.name,
                source_type='docx',
                page_number=None,
                additional_metadata=metadata
            )

            logger.info(f"Processed DOCX: {len(chunks)} chunks from {len(doc.paragraphs)} paragraphs")
            return chunks

        except Exception as e:
            logger.error(f"Error processing DOCX {file_path}: {str(e)}")
            raise

    def _create_chunks_from_text(
        self,
        text: str,
        source_file: str,
        source_type: str,
        page_number: Optional[int] = None,
        additional_metadata: Dict[str, Any] = None
    ) -> List[DocumentMetadata]:

        if not text.strip():
            return []

        chunks: List[DocumentMetadata] = []
        start = 0
        chunk_index = 0
        text_len = len(text)

        while start < text_len:
            end = min(start + self.chunk_size, text_len)
            if end < text_len:
                last_period = text.rfind('.', start, end)
                last_newline = text.rfind('\n', start, end)
                boundary = max(last_period, last_newline)
                if boundary > start + int(self.chunk_size * 0.5):
                    end = boundary + 1

            chunk_text = text[start:end].strip()

            if chunk_text:
                chunk_metadata = additional_metadata.copy() if additional_metadata else {}

                chunk = DocumentMetadata(
                    content=chunk_text,
                    source_file=source_file,
                    source_type=source_type,
                    page_number=page_number,
                    chunk_index=chunk_index,
                    start_char=start,
                    end_char=end - 1,
                    metadata=chunk_metadata
                )

                chunks.append(chunk)
                chunk_index += 1

            start = max(start + self.chunk_size - self.chunk_overlap, end)
            if start >= text_len:
                break

        return chunks

    def batch_process(self, file_paths: List[str]) -> List[DocumentMetadata]:
        all_chunks: List[DocumentMetadata] = []
        for file_path in file_paths:
            try:
                chunks = self.process_document(file_path)
                all_chunks.extend(chunks)
                logger.info(f"Successfully processed {file_path}: {len(chunks)} chunks")
            except Exception as e:
                logger.error(f"Failed to process {file_path}: {str(e)}")
                continue

        logger.info(f"Batch processing complete: {len(all_chunks)} total chunks from {len(file_paths)} files")
        return all_chunks


if __name__ == "__main__":
    processor = DocumentProcessor(chunk_size=800, chunk_overlap=100)

    try:
        chunks = processor.process_document("data/raft.pdf")
        sample_chunk = chunks[0]
        print(f"Sample chunk content: {sample_chunk.content[:200]}...")
        print(f"Citation info: {sample_chunk.get_citation_info()}")

    except Exception as e:
        print(f"Error: {e}")